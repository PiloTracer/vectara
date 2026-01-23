"""
Microsoft Office document extractors.
Supports DOCX, XLSX, and PPTX formats.
"""
import logging
from typing import List
from .base import BaseExtractor, ExtractedDocument

logger = logging.getLogger(__name__)


class OfficeExtractor(BaseExtractor):
    """
    Extract text from Microsoft Office documents (DOCX, XLSX, PPTX).
    """
    
    def supported_extensions(self) -> List[str]:
        return [
            ".docx", ".DOCX",
            ".xlsx", ".XLSX", 
            ".pptx", ".PPTX",
        ]
    
    async def extract(self, file_bytes: bytes, filename: str) -> ExtractedDocument:
        ext = self._get_extension(filename).lower()
        
        if ext == ".docx":
            return await self._extract_docx(file_bytes, filename)
        elif ext == ".xlsx":
            return await self._extract_xlsx(file_bytes, filename)
        elif ext == ".pptx":
            return await self._extract_pptx(file_bytes, filename)
        else:
            raise ValueError(f"Unsupported extension: {ext}")
    
    async def _extract_docx(self, file_bytes: bytes, filename: str) -> ExtractedDocument:
        """Extract text from Word documents."""
        try:
            from docx import Document
            import io
        except ImportError:
            raise ImportError("python-docx is required for DOCX extraction. Install with: pip install python-docx")
        
        doc = Document(io.BytesIO(file_bytes))
        paragraphs = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    paragraphs.append(" | ".join(row_text))
        
        content = "\n\n".join(paragraphs)
        
        return ExtractedDocument(
            content=content,
            metadata={"filename": filename, "type": "docx"},
            pages=None  # DOCX doesn't have strict page count without rendering
        )
    
    async def _extract_xlsx(self, file_bytes: bytes, filename: str) -> ExtractedDocument:
        """Extract text from Excel spreadsheets."""
        try:
            from openpyxl import load_workbook
            import io
        except ImportError:
            raise ImportError("openpyxl is required for XLSX extraction. Install with: pip install openpyxl")
        
        wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        sheet_texts = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = []
            
            for row in sheet.iter_rows(values_only=True):
                # Filter out None values and convert to strings
                row_values = [str(cell) if cell is not None else "" for cell in row]
                # Skip completely empty rows
                if any(v.strip() for v in row_values):
                    rows.append(" | ".join(row_values))
            
            if rows:
                sheet_text = f"## Sheet: {sheet_name}\n" + "\n".join(rows)
                sheet_texts.append(sheet_text)
        
        wb.close()
        content = "\n\n".join(sheet_texts)
        
        return ExtractedDocument(
            content=content,
            metadata={"filename": filename, "type": "xlsx", "sheets": len(wb.sheetnames)},
        )
    
    async def _extract_pptx(self, file_bytes: bytes, filename: str) -> ExtractedDocument:
        """Extract text from PowerPoint presentations."""
        try:
            from pptx import Presentation
            import io
        except ImportError:
            raise ImportError("python-pptx is required for PPTX extraction. Install with: pip install python-pptx")
        
        prs = Presentation(io.BytesIO(file_bytes))
        slide_texts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            
            if texts:
                slide_text = f"## Slide {slide_num}\n" + "\n".join(texts)
                slide_texts.append(slide_text)
        
        content = "\n\n".join(slide_texts)
        
        return ExtractedDocument(
            content=content,
            metadata={"filename": filename, "type": "pptx"},
            pages=len(prs.slides)
        )
